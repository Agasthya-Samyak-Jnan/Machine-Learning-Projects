from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# --- Helper to apply background color ---
def set_bg_color(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

# Slide 1
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
set_bg_color(slide1, 230, 240, 255)  # light blue background

title1 = slide1.shapes.title
body1 = slide1.placeholders[1]

title1.text = "Internship Overview"
title1.text_frame.paragraphs[0].font.size = Pt(44)
title1.text_frame.paragraphs[0].font.bold = True
title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 40, 90)

body1.text = (
    "Name: Agasthya Samyak Jnan\n"
    "USN: 02JSTUCS006\n"
    "Web Development Internship\n"
    "Duration: 08/09/2025 - 05/10/2025\n"
    "Internship Duration: 4 Weeks\n"
    "Organization: InternPe"
)

for p in body1.text_frame.paragraphs:
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(30, 30, 30)

# Slide 2
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
set_bg_color(slide2, 245, 235, 255)  # lavender background

title2 = slide2.shapes.title
body2 = slide2.placeholders[1]

title2.text = "Weekly Work Summary"
title2.text_frame.paragraphs[0].font.size = Pt(44)
title2.text_frame.paragraphs[0].font.bold = True
title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 0, 90)

body2.text = (
    "Week 1: Built a Calculator\n"
    "Week 2: Built an E-commerce Website"
)

for p in body2.text_frame.paragraphs:
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(50, 30, 70)

output_path = "/mnt/data/Styled_Internship_Presentation.pptx"
prs.save(output_path)

output_path
